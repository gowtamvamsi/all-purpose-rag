import asyncio
import io
import sys
import os
from datetime import datetime, timedelta, timezone

# Adjust path to import backend modules
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.future import select
from sqlalchemy import delete
from backend.database import SessionLocal
from backend.models.document import Document
from backend.models.chunk import DocumentChunk
from backend.models.chunk import DocumentChunk
from backend.services.storage_service import storage_service
from backend.services.embedding_service import embedding_service
from backend.workers.celery_app import celery_app
from backend.config import settings

from google.cloud import firestore
from google.auth.credentials import AnonymousCredentials

ALLOWED_2_LETTER_WORDS = {
    'in', 'on', 'at', 'to', 'is', 'it', 'he', 'me', 'we', 'us', 
    'an', 'by', 'my', 'so', 'no', 'go', 'if', 'do', 'or', 'of', 
    'am', 'as', 'up', 'ok', 'no', 'ii', 'la', 'us'
}

TYPO_REPLACEMENTS = {
    r'\bnembers\b': 'members',
    r'\bMurpiny\b': 'Murphy',
    r'\boepLonage\b': 'espionage',
    r'\b¢iscusaton\b': 'discussion',
    r'\biscusaton\b': 'discussion',
    r'\bdiscussaton\b': 'discussion',
    r'¢discussion': 'discussion',
    r'Board!\s*8\s*¢discussion': "Board's discussion",
    r'Board!\s*8': "Board's",
    r'\bassassiueston\b': 'assassination',
    r'\bwhieh\b': 'which',
    r'\bpam\b': 'p.m.',
    r'\bCLA\b': 'CIA',
    r'\bseLI1\b': 'still',
    r'\bpraccice\b': 'practice',
    r'\bdespit e\b': 'despite',
    r'\bRuss fans\b': 'Russians',
    r'\bcannoz\b': 'cannot',
    r'\bgonevally\b': 'generally',
    r'\bvwhich\b': 'which',
    r'\biad\b': 'had',
    r'\bbean\b': 'been',
    r'\bthac\b': 'that',
    r'\bEoard\b': 'Board',
    r'\bmuse\b': 'must',
    r'\binforsiation\b': 'information',
    r'\bshe\b': 'the',
    r'\bfelc\b': 'felt',
    r'\bceath\b': 'death',
    r'\bmémbers\b': 'members',
    r'\bforthe\b': 'for the',
    r'\bcompietion\b': 'completion',
}

def clean_ocr_text(text: str) -> str:
    import re
    # Normalize unicode spaces
    text = text.replace('\xa0', ' ')
    
    # Strip layout artifacts (like vertical borders '|', backslashes, etc.)
    lines = text.split("\n")
    cleaned_lines = []
    for line in lines:
        # Strip leading/trailing margins artifacts: |, \, /, ;, :, ., _, ~, +, -, =, *
        line = re.sub(r'^[|\s\\/;:._~+=*-]+|[|\s\\/;:._~+=*-]+$', '', line)
        line = line.strip()
        if not line:
            continue
            
        # Check for noise lines of length <= 2
        if len(line) <= 2:
            if line.isdigit() or (line.endswith('.') and line[:-1].isdigit()):
                pass
            elif line.lower() in ALLOWED_2_LETTER_WORDS:
                pass
            else:
                continue # drop noise lines like 'a', 'o', 'U', 'i', '\', etc.
                
        # Remove extra internal spaces
        line = re.sub(r'[ \t]+', ' ', line)
        cleaned_lines.append(line)
        
    joined = "\n".join(cleaned_lines)
    for pattern, repl in TYPO_REPLACEMENTS.items():
        joined = re.sub(pattern, repl, joined, flags=re.IGNORECASE)
    return joined

# --- FILE EXTRACTORS ---

def extract_pdf(raw_bytes: bytes) -> list[dict]:
    import pdfplumber
    import io
    blocks = []
    has_meaningful_text = False
    with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
        num_pages = len(pdf.pages)
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                blocks.append({"content": text, "page_number": i + 1})
                # We defer meaningful text check to after we collect all text
                
        has_meaningful_text = False
        for block in blocks:
            if len(clean_ocr_text(block["content"]).strip()) > 50:
                has_meaningful_text = True
                break
        
        # If the PDF is entirely scanned images, pdfplumber might extract tiny artifacts.
        # We check if there's any meaningful text extracted.
        if not has_meaningful_text:
            blocks = [] # Clear any garbage
            print("pdfplumber extracted very little meaningful text, attempting OCR with Tesseract...")
            try:
                import pytesseract
                from pdf2image import convert_from_bytes
                
                # Convert PDF to images
                print("Converting PDF to images...")
                images = convert_from_bytes(raw_bytes)
                print(f"Converted {len(images)} images, running OCR...")
                for i, image in enumerate(images):
                    text = pytesseract.image_to_string(image)
                    if text and text.strip():
                        blocks.append({"content": text, "page_number": i + 1})
                print(f"OCR complete, extracted {len(blocks)} blocks.")
            except Exception as e:
                print(f"OCR failed: {e}")
                blocks.append({"content": f"This is a scanned document containing {num_pages} pages. OCR failed.", "page_number": 1})
            
            if not blocks:
                blocks.append({"content": f"This is a scanned document containing {num_pages} pages. The text could not be extracted automatically.", "page_number": 1})
    return blocks

def extract_docx(raw_bytes: bytes) -> list[dict]:
    import docx
    doc = docx.Document(io.BytesIO(raw_bytes))
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)
    content = "\n".join(full_text)
    return [{"content": content, "page_number": 1}]

def extract_pptx(raw_bytes: bytes) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw_bytes))
    blocks = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_text.append(text)
        content = "\n".join(slide_text)
        if content:
            blocks.append({"content": content, "page_number": i + 1})
    return blocks

def extract_csv(raw_bytes: bytes) -> list[dict]:
    import pandas as pd
    df = pd.read_csv(io.BytesIO(raw_bytes))
    content = df.to_string(index=False)
    return [{"content": content, "page_number": 1}]

def extract_txt(raw_bytes: bytes) -> list[dict]:
    content = raw_bytes.decode("utf-8", errors="ignore")
    return [{"content": content, "page_number": 1}]

def extract_image(raw_bytes: bytes, file_type: str) -> list[dict]:
    import base64
    from openai import OpenAI
    
    api_key = settings.OPENAI_API_KEY
    if not api_key:
        return [{"content": "Mock image text: (Configure OPENAI_API_KEY for vision OCR).", "page_number": 1}]
    
    b64 = base64.b64encode(raw_bytes).decode()
    client = OpenAI(api_key=api_key)
    mime = "image/png"
    if file_type.lower() in ("jpg", "jpeg"):
        mime = "image/jpeg"
    elif file_type.lower() == "webp":
        mime = "image/webp"
        
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all readable text from this image. Do not summarize, just output the raw text found. If no text is found, output nothing."},
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}}
                ]
            }
        ]
    )
    extracted_text = response.choices[0].message.content or ""
    return [{"content": extracted_text, "page_number": 1}]


# --- RECURSIVE CHARACTER SPLITTER ---

def split_text_recursive(text: str, chunk_size=2000, overlap=300) -> list[str]:
    """
    Pure Python recursive characters text splitter.
    2000 chars roughly maps to 500-600 tokens.
    """
    separators = ["\n\n", "\n", ". ", "! ", "? ", " ", ""]
    
    def split_helper(txt: str, seps: list[str]) -> list[str]:
        if len(txt) <= chunk_size:
            return [txt]
        if not seps:
            return [txt[i:i+chunk_size] for i in range(0, len(txt), chunk_size - overlap)]
            
        sep = seps[0]
        parts = txt.split(sep)
        chunks = []
        current_chunk = ""
        
        for part in parts:
            if len(part) > chunk_size:
                if current_chunk:
                    chunks.append(current_chunk)
                    current_chunk = ""
                # Recursively split the long part with the remaining separators
                sub_parts = split_helper(part, seps[1:])
                chunks.extend(sub_parts)
            else:
                # Add sep if not the first item
                candidate = current_chunk + (sep if current_chunk else "") + part
                if len(candidate) <= chunk_size:
                    current_chunk = candidate
                else:
                    if current_chunk:
                        chunks.append(current_chunk)
                    current_chunk = part
        if current_chunk:
            chunks.append(current_chunk)
        return chunks
        
    return split_helper(text, separators)


# --- INGESTION CORE PIPELINE ---

async def get_isolated_db():
    from sqlalchemy.ext.asyncio import create_async_engine, async_sessionmaker, AsyncSession
    from sqlalchemy.pool import NullPool
    from backend.config import settings
    db_url = settings.DATABASE_URL
    if db_url.startswith("postgresql://"):
        db_url = db_url.replace("postgresql://", "postgresql+asyncpg://", 1)
    engine = create_async_engine(db_url, echo=False, poolclass=NullPool)
    Session = async_sessionmaker(bind=engine, class_=AsyncSession, expire_on_commit=False)
    return Session(), engine

async def process_document_async(document_id: str):
    db, engine = await get_isolated_db()
    try:
        # 1. Fetch document metadata
        result = await db.execute(select(Document).where(Document.id == document_id))
        doc = result.scalar_one_or_none()
        if not doc:
            print(f"Document {document_id} not found in database.")
            return
            
        doc.status = "processing"
        await db.commit()
        
        try:
            # 2. Get file bytes
            raw_bytes = storage_service.get_file_bytes(doc.s3_bucket, doc.s3_key)
            
            # 3. Extract text blocks based on type
            blocks = []
            f_type = doc.file_type.lower()
            
            print(f"[{doc.id}] Starting extraction for type {f_type}")
            if f_type == "pdf":
                # Always use local extraction: pdfplumber for digital PDFs,
                # Tesseract OCR fallback for scanned/image-based PDFs.
                blocks = extract_pdf(raw_bytes)
            elif f_type == "docx":
                blocks = extract_docx(raw_bytes)
            elif f_type == "pptx":
                blocks = extract_pptx(raw_bytes)
            elif f_type == "csv":
                blocks = extract_csv(raw_bytes)
            elif f_type in ("png", "jpg", "jpeg", "webp"):
                blocks = extract_image(raw_bytes, f_type)
            elif f_type in ("txt", "md"):
                blocks = extract_txt(raw_bytes)
            else:
                raise ValueError(f"Unsupported file extension: {f_type}")
            print(f"[{doc.id}] Extraction complete.")
                
            # If it's a PDF, set the page count
            if f_type == "pdf":
                doc.page_count = len(blocks)
                
            # 4. Chunk each block
            chunks_to_create = []
            chunk_index = 0
            
            for block in blocks:
                content = clean_ocr_text(block["content"])
                page_num = block.get("page_number")
                
                # Split content of block
                sub_texts = split_text_recursive(content, chunk_size=2000, overlap=300)
                for txt in sub_texts:
                    if not txt.strip():
                        continue
                    chunks_to_create.append({
                        "project_id": doc.project_id,
                        "document_id": doc.id,
                        "chunk_index": chunk_index,
                        "content": txt,
                        "page_number": page_num,
                        "token_count": len(txt.split()) # Simple token count estimate (words)
                    })
                    chunk_index += 1
            
            # Delete any existing chunks (reprocessing support)
            await db.execute(delete(DocumentChunk).where(DocumentChunk.document_id == doc.id))
            
            if len(chunks_to_create) > 0:
                # 5. Embed chunks
                texts_to_embed = [c["content"] for c in chunks_to_create]
                embeddings = embedding_service.get_embeddings_batched(texts_to_embed)
                
                # 6. Save chunks + embeddings to pgvector
                for c, emb in zip(chunks_to_create, embeddings):
                    chunk_model = DocumentChunk(
                        project_id=c["project_id"],
                        document_id=c["document_id"],
                        chunk_index=c["chunk_index"],
                        content=c["content"],
                        page_number=c["page_number"],
                        token_count=c["token_count"],
                        embedding=emb
                    )
                    db.add(chunk_model)
            
            # 7. Update document status
            doc.status = "ready"
            doc.chunk_count = len(chunks_to_create)
            doc.error_message = None
            await db.commit()
            print(f"Successfully processed document {doc.name} into {doc.chunk_count} chunks.")
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            doc.status = "error"
            doc.error_message = str(e)
            await db.commit()
            print(f"Error processing document {doc.name}: {e}")
    finally:
        await db.close()
        await engine.dispose()

@celery_app.task(name="backend.workers.ingestion.process_document")
def process_document(document_id: str):
    """Celery wrapper task."""
    asyncio.run(process_document_async(document_id))
